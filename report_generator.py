#!/usr/bin/env python3
"""
Professional Report Generator
专业级Review报告生成器

功能：
1. PPTX导出 - 包含封面、目录、Summary、详细检查项、匹配分析、逐Net详情、总结
2. PDF导出 - 同样完整的报告内容
3. 高质量图表和可视化
4. 符合工程文档标准
"""

import os
import io
from datetime import datetime
from typing import List, Dict, Tuple

import plotly.graph_objects as go
from plotly.subplots import make_subplots
import plotly.io as pio

# PPTX imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# PDF generation
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch, mm
    from reportlab.platypus import (
        SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer,
        PageBreak, Image, KeepTogether, ListFlowable, ListItem
    )
    from reportlab.pdfgen import canvas
    from reportlab.graphics.shapes import Drawing
    from reportlab.graphics.charts.barcharts import VerticalBarChart
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    print("Warning: reportlab not available, PDF export disabled")

from review_engine import (
    ProfessionalLayoutReviewEngine, ReviewSummary, Violation,
    MatchingAnalysis, NetRCData, Severity
)
from config_system import LayoutReviewConfig


# ============================================================================
# 颜色和样式定义
# ============================================================================

class ReportTheme:
    """报告主题颜色和字体"""

    # 字体配置 - 统一专业字体
    FONT_TITLE_CN = "Microsoft YaHei"      # 中文标题
    FONT_BODY_CN = "Microsoft YaHei"       # 中文字体
    FONT_TITLE_EN = "Times New Roman"       # 英文标题
    FONT_BODY_EN = "Times New Roman"        # 英文字体
    FONT_MONO = "Consolas"                  # 等宽字体 (用于数据/代码)

    # 主色调
    PRIMARY = RGBColor(0x2C, 0x3E, 0x50)      # 深蓝灰
    SECONDARY = RGBColor(0x34, 0x98, 0xDB)     # 蓝色
    ACCENT = RGBColor(0xE7, 0x4C, 0x3C)       # 红色
    SUCCESS = RGBColor(0x27, 0xAE, 0x60)       # 绿色
    WARNING = RGBColor(0xF3, 0x9C, 0x12)       # 橙色
    INFO = RGBColor(0x5D, 0xAD, 0xE2)          # 浅蓝
    
    # 中性色
    TEXT = RGBColor(0x2C, 0x3E, 0x50)          # 文本色
    TEXT_LIGHT = RGBColor(0x7F, 0x8C, 0x8D)    # 浅文本
    BACKGROUND = RGBColor(0xF8, 0xF9, 0xFA)    # 背景色
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)         # 白色
    BORDER = RGBColor(0xDE, 0xE2, 0xE6)        # 边框色
    
    # 严重程度颜色
    CRITICAL = RGBColor(0xC0, 0x39, 0x2B)      # 深红
    WARNING_COLOR = RGBColor(0xF3, 0x9C, 0x12) # 橙
    INFO_COLOR = RGBColor(0x34, 0x98, 0xDB)    # 蓝
    
    # ReportLab颜色 (PDF用) - 只在reportlab可用时定义
    if REPORTLAB_AVAILABLE:
        RL_CRITICAL = colors.HexColor('#C0392B')
        RL_WARNING = colors.HexColor('#F39C12')
        RL_INFO = colors.HexColor('#3498DB')
        RL_SUCCESS = colors.HexColor('#27AE60')
        RL_PRIMARY = colors.HexColor('#2C3E50')


# ============================================================================
# PPTX报告生成器
# ============================================================================

class PPTXReportGenerator:
    """PPTX格式报告生成器"""

    # 幻灯片尺寸 (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(self):
        self.prs = None
        self.theme = ReportTheme()

    def _set_font(self, paragraph, size, bold=False, color=None, name=None):
        """统一设置段落字体

        Args:
            paragraph: 段落对象
            size: 字体大小 (Pt)
            bold: 是否加粗
            color: 字体颜色 (RGBColor)
            name: 字体名称 (默认使用 TIMES NEW ROMAN)
        """
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        if color:
            paragraph.font.color.rgb = color
        # 使用 TIMES NEW ROMAN 作为默认英文字体
        paragraph.font.name = name or self.theme.FONT_TITLE_EN

    def create_report(self, engine: ProfessionalLayoutReviewEngine,
                      title: str = "Layout Review Report") -> 'PPTXReportGenerator':
        """创建完整报告"""
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT
        
        summary = engine._generate_summary()
        
        # 1. 封面
        self._add_cover_slide(title, engine.config)
        
        # 2. 目录
        self._add_toc_slide()
        
        # 3. Executive Summary
        self._add_summary_slide(summary, engine)
        
        # 4. Net Statistics Overview
        self._add_net_statistics_slide(engine)
        
        # 5. Violation Overview
        if summary.total_violations > 0:
            self._add_violation_summary_slide(summary, engine)
        
        # 6. Critical Issues (详细)
        critical_violations = [v for v in engine.violations if v.severity == Severity.CRITICAL]
        if critical_violations:
            self._add_violation_detail_slides(critical_violations, "Critical Issues")
        
        # 7. Warnings (详细)
        warning_violations = [v for v in engine.violations if v.severity == Severity.WARNING]
        if warning_violations:
            self._add_violation_detail_slides(warning_violations[:10], "Warning Issues")  # 最多10个
        
        # 8. Matching Analysis
        if engine.matching_results:
            self._add_matching_summary_slide(engine.matching_results)
            self._add_matching_detail_slides(engine.matching_results[:6])  # 最多6个

        # 9. Net Details - per-net detailed pages
        self._add_net_detail_slides(engine)

        # 10. Recommendations
        self._add_recommendations_slide(engine)

        # 11. 结束页
        self._add_final_slide()
        
        return self
    
    def _add_cover_slide(self, title: str, config: LayoutReviewConfig):
        """添加封面"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank

        # 纯色背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme.PRIMARY
        bg.line.fill.background()

        # 主标题 - 统一字体
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.2),
            Inches(12), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.theme.WHITE
        p.font.name = self.theme.FONT_TITLE_EN
        p.alignment = PP_ALIGN.LEFT

        # 副标题
        sub_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.5),
            Inches(12), Inches(0.8)
        )
        tf = sub_box.text_frame
        tf.text = f"{config.tech_config.name} | {config.tech_config.node} Process"
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = self.theme.SECONDARY
        p.font.name = self.theme.FONT_BODY_EN
        p.alignment = PP_ALIGN.LEFT

        # 底部信息
        info_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(6.2),
            Inches(12), Inches(0.8)
        )
        tf = info_box.text_frame
        tf.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = self.theme.TEXT_LIGHT
        p.font.name = self.theme.FONT_BODY_EN
        p.alignment = PP_ALIGN.LEFT
    
    def _add_toc_slide(self):
        """添加目录页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.4),
            Inches(12), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = "Table of Contents"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.theme.PRIMARY
        p.font.name = self.theme.FONT_TITLE_EN

        # 目录项
        contents = [
            ("01", "Executive Summary", "High-level review overview"),
            ("02", "Net Statistics", "RC analysis and layer usage"),
            ("03", "Violation Analysis", "Critical and warning issues"),
            ("04", "Matching Analysis", "Signal pair matching results"),
            ("05", "Net Details", "Per-net detailed analysis with visualizations"),
            ("06", "Recommendations", "Action items and next steps"),
        ]

        y_pos = Inches(1.4)
        for num, title, desc in contents:
            # 编号
            num_box = slide.shapes.add_textbox(
                Inches(0.7), y_pos,
                Inches(0.7), Inches(0.4)
            )
            tf = num_box.text_frame
            tf.text = num
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.theme.SECONDARY
            p.font.name = self.theme.FONT_TITLE_EN

            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(1.6), y_pos,
                Inches(5), Inches(0.4)
            )
            tf = title_box.text_frame
            tf.text = title
            p = tf.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.theme.TEXT
            p.font.name = self.theme.FONT_TITLE_EN

            # 描述
            desc_box = slide.shapes.add_textbox(
                Inches(1.6), y_pos + Inches(0.35),
                Inches(8), Inches(0.3)
            )
            tf = desc_box.text_frame
            tf.text = desc
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = self.theme.TEXT_LIGHT
            p.font.name = self.theme.FONT_BODY_EN

            y_pos += Inches(0.9)
    
    def _add_summary_slide(self, summary: ReviewSummary, engine: ProfessionalLayoutReviewEngine):
        """添加Executive Summary"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题栏
        self._add_slide_header(slide, "Executive Summary")
        
        # 统计卡片
        stats = [
            ("Total Nets", str(summary.total_nets), self.theme.PRIMARY),
            ("Critical", str(summary.critical_count), self.theme.CRITICAL if summary.critical_count > 0 else self.theme.SUCCESS),
            ("Warnings", str(summary.warning_count), self.theme.WARNING if summary.warning_count > 0 else self.theme.SUCCESS),
            ("Info", str(summary.info_count), self.theme.INFO_COLOR),
        ]
        
        card_width = Inches(2.8)
        card_height = Inches(1.3)
        start_left = Inches(0.5)
        start_top = Inches(1.2)
        
        for i, (label, value, color) in enumerate(stats):
            left = start_left + i * (card_width + Inches(0.2))
            
            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, start_top,
                card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.theme.WHITE
            card.line.color.rgb = color
            card.line.width = Pt(2)
            
            # 数值
            value_box = slide.shapes.add_textbox(
                left, start_top + Inches(0.15),
                card_width, Inches(0.7)
            )
            tf = value_box.text_frame
            tf.text = value
            p = tf.paragraphs[0]
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER
            
            # 标签
            label_box = slide.shapes.add_textbox(
                left, start_top + Inches(0.85),
                card_width, Inches(0.4)
            )
            tf = label_box.text_frame
            tf.text = label
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme.TEXT
            p.alignment = PP_ALIGN.CENTER
        
        # RC统计
        rc_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(2.8),
            Inches(6), Inches(2.2)
        )
        rc_box.fill.solid()
        rc_box.fill.fore_color.rgb = self.theme.BACKGROUND
        rc_box.line.fill.background()
        
        rc_title = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.0),
            Inches(5.5), Inches(0.4)
        )
        tf = rc_title.text_frame
        tf.text = "RC Analysis Summary"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.theme.PRIMARY
        
        rc_text = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.4),
            Inches(5.5), Inches(1.5)
        )
        tf = rc_text.text_frame
        tf.text = f"Resistance Range: {summary.total_resistance_range[0]:.1f} - {summary.total_resistance_range[1]:.1f} Ω\n"
        tf.text += f"Average Resistance: {summary.avg_resistance:.1f} Ω\n\n"
        tf.text += f"Capacitance Range: {summary.total_capacitance_range[0]:.1f} - {summary.total_capacitance_range[1]:.1f} fF\n"
        tf.text += f"Average Capacitance: {summary.avg_capacitance:.1f} fF\n\n"
        tf.text += f"Matching Pairs: {summary.matching_pairs_analyzed}"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = self.theme.TEXT
        
        # 违规分类饼图区域
        if summary.violations_by_type:
            chart_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.8), Inches(2.8),
                Inches(6), Inches(2.2)
            )
            chart_box.fill.solid()
            chart_box.fill.fore_color.rgb = self.theme.BACKGROUND
            chart_box.line.fill.background()
            
            chart_title = slide.shapes.add_textbox(
                Inches(7.0), Inches(3.0),
                Inches(5.5), Inches(0.4)
            )
            tf = chart_title.text_frame
            tf.text = "Violations by Type"
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme.PRIMARY
            
            # 简单的文字列表代替饼图
            y_pos = Inches(3.4)
            for vtype, count in sorted(summary.violations_by_type.items(), key=lambda x: -x[1])[:4]:
                type_box = slide.shapes.add_textbox(
                    Inches(7.0), y_pos,
                    Inches(5), Inches(0.25)
                )
                tf = type_box.text_frame
                tf.text = f"- {vtype}: {count}"
                p = tf.paragraphs[0]
                p.font.size = Pt(11)
                p.font.color.rgb = self.theme.TEXT
                y_pos += Inches(0.3)
        
        # 状态总结
        status_text = "✓ PASSED" if summary.critical_count == 0 and summary.warning_count == 0 else \
                     "⚠ REVIEW REQUIRED" if summary.critical_count == 0 else \
                     "✗ CRITICAL ISSUES FOUND"
        status_color = self.theme.SUCCESS if summary.critical_count == 0 and summary.warning_count == 0 else \
                      self.theme.WARNING if summary.critical_count == 0 else self.theme.CRITICAL
        
        status_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.2),
            Inches(12), Inches(0.6)
        )
        tf = status_box.text_frame
        tf.text = f"Overall Status: {status_text}"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = status_color
    
    def _add_net_statistics_slide(self, engine: ProfessionalLayoutReviewEngine):
        """添加Net统计表"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_header(slide, "Net Statistics Overview")

        # 获取数据
        table_data = engine.get_net_info_table()
        if not table_data:
            return

        # 显示前15个net
        display_data = table_data[:15]

        # 创建表格 - 扩展列包含时序数据
        headers = ['Net Name', 'R (Ω)', 'C (fF)', 'Length (μm)', 'Via Count',
                   'tau_rc (ps)', 't_pd (ps)', 't_rise (ps)', 't_fall (ps)', 'Violations']
        rows = len(display_data) + 1
        cols = len(headers)

        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.4), Inches(1.3),
            Inches(12.5), Inches(5.8)
        ).table

        # 设置列宽 - 调整为10列
        col_widths = [Inches(2.3), Inches(1.0), Inches(1.0), Inches(1.2), Inches(0.9),
                     Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.1)]
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

        # 表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.color.rgb = self.theme.WHITE
            cell.text_frame.paragraphs[0].font.name = self.theme.FONT_TITLE_EN
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.PRIMARY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 数据行
        for row_idx, row_data in enumerate(display_data, 1):
            values = [
                row_data['Net Name'],
                f"{row_data['Total R (Ω)']:.2f}",
                f"{row_data['Total C (fF)']:.1f}",
                f"{row_data['Length (μm)']:.1f}",
                str(row_data['Via Count']),
                f"{row_data.get('tau_rc (ps)', 0):.4f}",
                f"{row_data.get('tpd_50% (ps)', 0):.4f}",
                f"{row_data.get('trise (ps)', 0):.4f}",
                f"{row_data.get('tfall (ps)', 0):.4f}",
                str(row_data['Violations'])
            ]

            for col_idx, value in enumerate(values):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                cell.text_frame.paragraphs[0].font.size = Pt(8)
                cell.text_frame.paragraphs[0].font.name = self.theme.FONT_BODY_EN
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 斑马纹
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.theme.BACKGROUND

                # 违规数高亮
                if col_idx == 9 and row_data['Violations'] > 0:
                    cell.text_frame.paragraphs[0].font.color.rgb = self.theme.CRITICAL
                    cell.text_frame.paragraphs[0].font.bold = True

        # 如果还有更多
        if len(table_data) > 15:
            more_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(7.2),
                Inches(12), Inches(0.3)
            )
            tf = more_box.text_frame
            tf.text = f"... and {len(table_data) - 15} more nets"
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = self.theme.TEXT_LIGHT
    
    def _add_violation_summary_slide(self, summary: ReviewSummary, engine: ProfessionalLayoutReviewEngine):
        """添加违规汇总"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_header(slide, "Violation Overview")
        
        # 严重程度分布
        severity_data = [
            ("Critical", summary.critical_count, self.theme.CRITICAL),
            ("Warning", summary.warning_count, self.theme.WARNING_COLOR),
            ("Info", summary.info_count, self.theme.INFO_COLOR),
        ]
        
        y_pos = Inches(1.5)
        for severity, count, color in severity_data:
            # 颜色块
            block = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), y_pos,
                Inches(0.3), Inches(0.3)
            )
            block.fill.solid()
            block.fill.fore_color.rgb = color
            block.line.fill.background()
            
            # 文本
            text_box = slide.shapes.add_textbox(
                Inches(1.0), y_pos - Inches(0.05),
                Inches(4), Inches(0.4)
            )
            tf = text_box.text_frame
            tf.text = f"{severity}: {count}"
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = self.theme.TEXT
            
            y_pos += Inches(0.5)
        
        # 按类型排序的违规
        type_box = slide.shapes.add_textbox(
            Inches(6), Inches(1.5),
            Inches(6.5), Inches(5.5)
        )
        tf = type_box.text_frame
        tf.text = "Top Violation Types\n"
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.theme.PRIMARY
        
        for vtype, count in sorted(summary.violations_by_type.items(), key=lambda x: -x[1])[:10]:
            p = tf.add_paragraph()
            p.text = f"- {vtype}: {count}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.theme.TEXT
    
    def _add_violation_detail_slides(self, violations: List[Violation], title_prefix: str):
        """添加违规详情页"""
        # 每页显示3个
        for i in range(0, len(violations), 3):
            batch = violations[i:i+3]
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self._add_slide_header(slide, f"{title_prefix} ({i+1}-{min(i+3, len(violations))})")
            
            y_pos = Inches(1.3)
            for v in batch:
                # 违规卡片
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.5), y_pos,
                    Inches(12.3), Inches(1.8)
                )
                
                card.fill.solid()
                
                if v.severity == Severity.CRITICAL:
                    card.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
                    border_color = self.theme.CRITICAL
                elif v.severity == Severity.WARNING:
                    card.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
                    border_color = self.theme.WARNING_COLOR
                else:
                    card.fill.fore_color.rgb = self.theme.BACKGROUND
                    border_color = self.theme.INFO_COLOR
                card.line.color.rgb = border_color
                card.line.width = Pt(1.5)
                
                # 标题行
                title_box = slide.shapes.add_textbox(
                    Inches(0.7), y_pos + Inches(0.1),
                    Inches(11.5), Inches(0.35)
                )
                tf = title_box.text_frame
                tf.text = f"[{v.rule_id}] {v.rule_name}"
                p = tf.paragraphs[0]
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = border_color
                
                # Net名称
                net_box = slide.shapes.add_textbox(
                    Inches(0.7), y_pos + Inches(0.45),
                    Inches(3), Inches(0.3)
                )
                tf = net_box.text_frame
                tf.text = f"Net: {v.net_name}"
                p = tf.paragraphs[0]
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = self.theme.TEXT
                
                # 消息
                msg_box = slide.shapes.add_textbox(
                    Inches(0.7), y_pos + Inches(0.75),
                    Inches(11.5), Inches(0.5)
                )
                tf = msg_box.text_frame
                tf.text = v.message
                p = tf.paragraphs[0]
                p.font.size = Pt(10)
                p.font.color.rgb = self.theme.TEXT
                
                # 建议
                if v.suggestion:
                    sug_box = slide.shapes.add_textbox(
                        Inches(0.7), y_pos + Inches(1.2),
                        Inches(11.5), Inches(0.4)
                    )
                    tf = sug_box.text_frame
                    tf.text = f"[Suggestion] {v.suggestion}"
                    p = tf.paragraphs[0]
                    p.font.size = Pt(9)
                    p.font.italic = True
                    p.font.color.rgb = self.theme.SUCCESS
                
                y_pos += Inches(2.0)
    
    def _add_matching_summary_slide(self, matching_results: List[MatchingAnalysis]):
        """添加匹配分析汇总"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_header(slide, "Matching Analysis Summary")
        
        # 统计
        scores = [m.match_score for m in matching_results]
        avg_score = sum(scores) / len(scores) if scores else 0
        poor_count = len([s for s in scores if s < 70])
        
        # 分数分布
        score_ranges = [
            ("Excellent (90-100)", len([s for s in scores if 90 <= s <= 100]), self.theme.SUCCESS),
            ("Good (70-89)", len([s for s in scores if 70 <= s < 90]), self.theme.INFO_COLOR),
            ("Fair (50-69)", len([s for s in scores if 50 <= s < 70]), self.theme.WARNING_COLOR),
            ("Poor (<50)", len([s for s in scores if s < 50]), self.theme.CRITICAL),
        ]
        
        y_pos = Inches(1.5)
        for label, count, color in score_ranges:
            block = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), y_pos,
                Inches(0.3), Inches(0.3)
            )
            block.fill.solid()
            block.fill.fore_color.rgb = color
            block.line.fill.background()
            
            text_box = slide.shapes.add_textbox(
                Inches(1.0), y_pos - Inches(0.05),
                Inches(5), Inches(0.4)
            )
            tf = text_box.text_frame
            tf.text = f"{label}: {count} pairs"
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme.TEXT
            
            y_pos += Inches(0.45)
        
        # 平均分
        avg_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.5),
            Inches(5), Inches(0.8)
        )
        tf = avg_box.text_frame
        tf.text = f"Average Match Score: {avg_score:.1f}"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.theme.PRIMARY if avg_score >= 80 else self.theme.WARNING
    
    def _add_matching_detail_slides(self, matching_results: List[MatchingAnalysis]):
        """添加匹配详情"""
        for match in matching_results:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            
            score_color = self.theme.SUCCESS if match.match_score >= 80 else \
                         self.theme.WARNING if match.match_score >= 60 else self.theme.CRITICAL
            
            self._add_slide_header(slide, f"Matching: {match.net1} ↔ {match.net2}")
            
            # 大分数显示
            score_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3),
                Inches(3), Inches(1.2)
            )
            tf = score_box.text_frame
            tf.text = f"{match.match_score:.0f}"
            p = tf.paragraphs[0]
            p.font.size = Pt(60)
            p.font.bold = True
            p.font.color.rgb = score_color
            
            label_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.3),
                Inches(3), Inches(0.4)
            )
            tf = label_box.text_frame
            tf.text = "Match Score"
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme.TEXT_LIGHT
            
            # 详细指标
            metrics = [
                ("Length Ratio", f"{match.length_ratio:.3f}"),
                ("Resistance Ratio", f"{match.resistance_ratio:.3f}"),
                ("Capacitance Ratio", f"{match.capacitance_ratio:.3f}"),
                ("Via Count Diff", str(match.via_count_diff)),
                ("Centroid Distance", f"{match.centroid_distance:.2f} μm"),
                ("BBox Similarity", f"{match.bbox_similarity:.2f}"),
            ]
            
            y_pos = Inches(1.4)
            for metric, value in metrics:
                metric_box = slide.shapes.add_textbox(
                    Inches(4), y_pos,
                    Inches(3), Inches(0.25)
                )
                tf = metric_box.text_frame
                tf.text = metric
                p = tf.paragraphs[0]
                p.font.size = Pt(10)
                p.font.color.rgb = self.theme.TEXT_LIGHT
                
                value_box = slide.shapes.add_textbox(
                    Inches(7), y_pos,
                    Inches(3), Inches(0.25)
                )
                tf = value_box.text_frame
                tf.text = value
                p = tf.paragraphs[0]
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = self.theme.TEXT
                
                y_pos += Inches(0.35)
            
            # 问题和建议
            if match.issues:
                issue_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(4.0),
                    Inches(12), Inches(0.4)
                )
                tf = issue_box.text_frame
                tf.text = "Issues Found:"
                p = tf.paragraphs[0]
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = self.theme.CRITICAL
                
                y_pos = Inches(4.4)
                for issue in match.issues[:3]:  # 最多3个
                    issue_text = slide.shapes.add_textbox(
                        Inches(0.7), y_pos,
                        Inches(11.5), Inches(0.25)
                    )
                    tf = issue_text.text_frame
                    tf.text = f"- {issue}"
                    p = tf.paragraphs[0]
                    p.font.size = Pt(10)
                    p.font.color.rgb = self.theme.TEXT
                    y_pos += Inches(0.3)
    
    def _add_rc_distribution_slide(self, engine: ProfessionalLayoutReviewEngine):
        """添加RC分布图表"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_header(slide, "RC Distribution Analysis")
        
        # 使用Plotly生成图表并转换为图片
        resistances = [rc.total_resistance for rc in engine.net_rc_data.values()]
        capacitances = [rc.total_capacitance for rc in engine.net_rc_data.values()]
        
        if resistances and capacitances:
            fig = make_subplots(rows=1, cols=2, 
                               subplot_titles=('Resistance Distribution', 'Capacitance Distribution'))
            
            # 电阻直方图
            fig.add_trace(go.Histogram(
                x=resistances,
                nbinsx=20,
                name='Resistance (Ω)',
                marker_color='#3498DB'
            ), row=1, col=1)
            
            # 电容直方图
            fig.add_trace(go.Histogram(
                x=capacitances,
                nbinsx=20,
                name='Capacitance (fF)',
                marker_color='#27AE60'
            ), row=1, col=2)
            
            fig.update_layout(
                showlegend=False,
                height=500,
                width=1100,
                plot_bgcolor='white',
                paper_bgcolor='white'
            )
            
            # 尝试导出图片
            try:
                img_bytes = pio.to_image(fig, format='png', scale=2)
                img_stream = io.BytesIO(img_bytes)
                slide.shapes.add_picture(
                    img_stream,
                    Inches(0.5), Inches(1.3),
                    width=Inches(12.3)
                )
            except Exception as e:
                # 如果导出失败，显示文字
                text_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(3),
                    Inches(12), Inches(1)
                )
                tf = text_box.text_frame
                tf.text = f"Chart generation requires kaleido\nError: {str(e)}"
                p = tf.paragraphs[0]
                p.font.size = Pt(14)
                p.font.color.rgb = self.theme.TEXT_LIGHT
    
    def _add_net_ranking_slide(self, engine: ProfessionalLayoutReviewEngine):
        """添加Net排名"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_header(slide, "Net Rankings - Top Issues")
        
        # 按违规数排序
        net_violations = [(net, len([v for v in engine.violations if v.net_name == net])) 
                         for net in engine.nets]
        net_violations.sort(key=lambda x: -x[1])
        
        if net_violations:
            # 标题
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.4),
                Inches(6), Inches(0.4)
            )
            tf = title_box.text_frame
            tf.text = "Nets by Violation Count"
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme.PRIMARY
            
            # 表格
            rows = min(11, len(net_violations) + 1)
            table = slide.shapes.add_table(
                rows, 3,
                Inches(0.5), Inches(1.9),
                Inches(5.5), Inches(5)
            ).table
            
            # 表头
            headers = ['Rank', 'Net Name', 'Violations']
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = h
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.color.rgb = self.theme.WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.theme.PRIMARY
            
            # 数据
            for rank, (net, count) in enumerate(net_violations[:10], 1):
                table.cell(rank, 0).text = str(rank)
                table.cell(rank, 1).text = net
                table.cell(rank, 2).text = str(count)
                
                for col in range(3):
                    cell = table.cell(rank, col)
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    if count > 0:
                        cell.text_frame.paragraphs[0].font.color.rgb = self.theme.CRITICAL
        
        # 按电阻排序
        net_resistances = [(net, rc.total_resistance) 
                          for net, rc in engine.net_rc_data.items()]
        net_resistances.sort(key=lambda x: -x[1])
        
        if net_resistances:
            title_box = slide.shapes.add_textbox(
                Inches(6.5), Inches(1.4),
                Inches(6), Inches(0.4)
            )
            tf = title_box.text_frame
            tf.text = "Nets by Resistance"
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.theme.PRIMARY
            
            # 表格
            rows = min(11, len(net_resistances) + 1)
            table = slide.shapes.add_table(
                rows, 3,
                Inches(6.5), Inches(1.9),
                Inches(5.5), Inches(5)
            ).table
            
            headers = ['Rank', 'Net Name', 'R (Ω)']
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = h
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.color.rgb = self.theme.WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.theme.PRIMARY
            
            for rank, (net, r) in enumerate(net_resistances[:10], 1):
                table.cell(rank, 0).text = str(rank)
                table.cell(rank, 1).text = net
                table.cell(rank, 2).text = f"{r:.1f}"
                
                for col in range(3):
                    cell = table.cell(rank, col)
                    cell.text_frame.paragraphs[0].font.size = Pt(9)

    def _add_net_detail_slides(self, engine: ProfessionalLayoutReviewEngine):
        """为每个net添加详细的分析页面

        每页包含:
        - 标题: Net: XXXXX
        - 三张截图: 全貌、密度最大区域、违例区域
        - 描述: Layer/shapes信息、RC结果
        - 结论: Critical/Warning/Info数量
        """
        for net_name in sorted(engine.nets.keys()):
            polygons = engine.nets.get(net_name, [])
            rc_data = engine.net_rc_data.get(net_name)
            net_violations = [v for v in engine.violations if v.net_name == net_name]

            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

            # 标题: Net: XXXXX
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(12), Inches(0.6)
            )
            tf = title_box.text_frame
            tf.text = f"Net: {net_name}"
            p = tf.paragraphs[0]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.theme.PRIMARY

            # 统计违例数量
            critical_count = len([v for v in net_violations if v.severity == Severity.CRITICAL])
            warning_count = len([v for v in net_violations if v.severity == Severity.WARNING])
            info_count = len([v for v in net_violations if v.severity == Severity.INFO])

            # 在右上角显示违例统计
            stats_box = slide.shapes.add_textbox(
                Inches(10), Inches(0.35),
                Inches(3), Inches(0.5)
            )
            tf = stats_box.text_frame
            stats_text = f"Critical: {critical_count} | Warning: {warning_count} | Info: {info_count}"
            tf.text = stats_text
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            if critical_count > 0:
                p.font.color.rgb = self.theme.CRITICAL
            elif warning_count > 0:
                p.font.color.rgb = self.theme.WARNING_COLOR
            else:
                p.font.color.rgb = self.theme.SUCCESS
            p.alignment = PP_ALIGN.RIGHT

            # 生成可视化图片
            try:
                from core.report_visualization import create_polygons_figure
                import plotly.io as pio

                # 获取该net的多边形按层分组
                layer_polys: Dict[str, List] = {}
                for poly in polygons:
                    if poly.layer not in layer_polys:
                        layer_polys[poly.layer] = []
                    layer_polys[poly.layer].append(poly)

                if not layer_polys:
                    continue

                # 计算整体bbox
                all_bboxes = [p.bbox for p in polygons]
                xmin = min(b[0] for b in all_bboxes)
                ymin = min(b[1] for b in all_bboxes)
                xmax = max(b[2] for b in all_bboxes)
                ymax = max(b[3] for b in all_bboxes)

                margin = 0.1 * max(xmax - xmin, ymax - ymin) if xmax > xmin and ymax > ymin else 1.0

                # 1. 全貌截图 (Fit view)
                fig_fit = create_polygons_figure(
                    polygons,
                    title=f"{net_name} - Full View",
                    x_range=[xmin - margin, xmax + margin],
                    y_range=[ymin - margin, ymax + margin],
                    showgrid=False
                )
                img_fit = pio.to_image(fig_fit, format='png', scale=1)

                # 2. 密度最大区域截图
                # 找到shape最密集的区域
                densest_layer = max(layer_polys.keys(), key=lambda l: len(layer_polys[l]))
                densest_polys = layer_polys[densest_layer]

                # 计算密集区域bbox
                d_bboxes = [p.bbox for p in densest_polys]
                dxmin = min(b[0] for b in d_bboxes)
                dymin = min(b[1] for b in d_bboxes)
                dxmax = max(b[2] for b in d_bboxes)
                dymax = max(b[3] for b in d_bboxes)

                fig_dense = create_polygons_figure(
                    densest_polys,
                    title=f"Highest Density Area - {densest_layer}",
                    x_range=[dxmin - margin * 0.5, dxmax + margin * 0.5],
                    y_range=[dymin - margin * 0.5, dymax + margin * 0.5],
                    showgrid=False
                )
                img_dense = pio.to_image(fig_dense, format='png', scale=1)

                # 3. 违例区域截图
                violation_polys = []
                for v in net_violations:
                    violation_polys.extend(v.polygons)

                if violation_polys:
                    v_bboxes = [p.bbox for p in violation_polys]
                    vxmin = min(b[0] for b in v_bboxes)
                    vymin = min(b[1] for b in v_bboxes)
                    vxmax = max(b[2] for b in v_bboxes)
                    vymax = max(b[3] for b in v_bboxes)

                    fig_viol = create_polygons_figure(
                        violation_polys,
                        title=f"Violation Areas ({len(violation_polys)} shapes)",
                        x_range=[vxmin - margin * 0.5, vxmax + margin * 0.5],
                        y_range=[vymin - margin * 0.5, vymax + margin * 0.5],
                        showgrid=False
                    )
                    img_viol = pio.to_image(fig_viol, format='png', scale=1)
                else:
                    # 没有违例时显示全貌缩略图
                    img_viol = img_fit

                # 在幻灯片上放置三张图片 (同一行)
                img_height = Inches(3.2)
                img_width = Inches(4.0)

                # 图片1: 全貌
                stream1 = io.BytesIO(img_fit)
                slide.shapes.add_picture(
                    stream1,
                    Inches(0.4), Inches(1.1),
                    width=img_width, height=img_height
                )

                # 图片2: 密度最大区域
                stream2 = io.BytesIO(img_dense)
                slide.shapes.add_picture(
                    stream2,
                    Inches(4.6), Inches(1.1),
                    width=img_width, height=img_height
                )

                # 图片3: 违例区域
                stream3 = io.BytesIO(img_viol)
                slide.shapes.add_picture(
                    stream3,
                    Inches(8.8), Inches(1.1),
                    width=img_width, height=img_height
                )

                # 图片标签
                labels = ["Full View (Fit)", f"Dense Area ({densest_layer})", "Violation Areas"]
                for i, label in enumerate(labels):
                    label_box = slide.shapes.add_textbox(
                        Inches(0.4 + i * 4.2), Inches(4.4),
                        Inches(4.0), Inches(0.3)
                    )
                    tf = label_box.text_frame
                    tf.text = label
                    p = tf.paragraphs[0]
                    p.font.size = Pt(9)
                    p.font.bold = True
                    p.font.color.rgb = self.theme.TEXT
                    p.alignment = PP_ALIGN.CENTER

            except Exception as e:
                # 可视化失败时显示占位符
                placeholder = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.5),
                    Inches(12), Inches(3)
                )
                tf = placeholder.text_frame
                tf.text = f"Visualization generation requires kaleido\nError: {str(e)[:100]}"
                p = tf.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = self.theme.TEXT_LIGHT

            # 描述区域: Layer/shapes信息和RC结果
            desc_y = Inches(4.8)

            # Left column: Layer info
            layer_info_box = slide.shapes.add_textbox(
                Inches(0.5), desc_y,
                Inches(4), Inches(1.8)
            )
            tf = layer_info_box.text_frame
            tf.text = "Layer Information\n"
            tf.text += "-" * 25 + "\n"

            for layer, polys_on_layer in layer_polys.items():
                tf.text += f"{layer}: {len(polys_on_layer)} shapes\n"
            tf.text += f"\nTotal: {len(polygons)} shapes"

            for p_obj in tf.paragraphs:
                p_obj.font.size = Pt(10)
                p_obj.font.color.rgb = self.theme.TEXT
                p_obj.font.name = "Consolas"

            # Middle column: RC results
            if rc_data:
                rc_info_box = slide.shapes.add_textbox(
                    Inches(4.7), desc_y,
                    Inches(4), Inches(2.2)
                )
                tf = rc_info_box.text_frame
                tf.text = "RC Results\n"
                tf.text += "-" * 25 + "\n"
                tf.text += f"R Total: {rc_data.total_resistance:.2f} Ω\n"
                tf.text += f"C Total: {rc_data.total_capacitance:.2f} fF\n"
                tf.text += f"Length: {rc_data.total_length:.2f} μm\n"
                tf.text += f"Area: {rc_data.total_area:.2f} μm²\n"
                tf.text += f"Vias: {rc_data.via_count}\n"
                tf.text += f"\nTiming Analysis\n"
                tf.text += f"  τ_rc: {rc_data.tau_rc:.3f} ps\n"
                tf.text += f"  t_pd@50%: {rc_data.tpd_50:.3f} ps"

                for p_obj in tf.paragraphs:
                    p_obj.font.size = Pt(10)
                    p_obj.font.color.rgb = self.theme.TEXT
                    p_obj.font.name = "Consolas"

            # Right column: Violations summary
            viol_info_box = slide.shapes.add_textbox(
                Inches(8.9), desc_y,
                Inches(4), Inches(1.8)
            )
            tf = viol_info_box.text_frame
            tf.text = "Violation Summary\n"
            tf.text += "-" * 25 + "\n"
            tf.text += f"Critical: {critical_count}\n"
            tf.text += f"Warning: {warning_count}\n"
            tf.text += f"Info: {info_count}"

            if net_violations:
                tf.text += "\n\nDetails:\n"
                for v in net_violations[:5]:
                    short_msg = v.message[:40] + "..." if len(v.message) > 40 else v.message
                    tf.text += f"[{v.severity.value.upper()}] {short_msg}\n"

            for p_obj in tf.paragraphs:
                p_obj.font.size = Pt(10)
                p_obj.font.color.rgb = self.theme.TEXT
                p_obj.font.name = "Consolas"

    def _add_recommendations_slide(self, engine: ProfessionalLayoutReviewEngine):
        """添加建议页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_header(slide, "Recommendations & Action Items")
        
        # 收集所有建议
        all_suggestions = []
        for v in engine.violations:
            if v.suggestion and v.suggestion not in all_suggestions:
                all_suggestions.append((v.severity, v.suggestion))
        
        for m in engine.matching_results:
            for sug in m.suggestions:
                if sug not in [s[1] for s in all_suggestions]:
                    all_suggestions.append((Severity.WARNING, sug))
        
        # 按严重程度排序
        severity_order = {Severity.CRITICAL: 0, Severity.WARNING: 1, Severity.INFO: 2}
        all_suggestions.sort(key=lambda x: severity_order.get(x[0], 3))
        
        y_pos = Inches(1.5)
        for severity, suggestion in all_suggestions[:8]:  # 最多8个
            color = self.theme.CRITICAL if severity == Severity.CRITICAL else \
                   self.theme.WARNING_COLOR if severity == Severity.WARNING else self.theme.INFO_COLOR
            
            # 项目符号
            bullet = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), y_pos + Inches(0.08),
                Inches(0.15), Inches(0.15)
            )
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = color
            bullet.line.fill.background()
            
            # 建议文本
            text_box = slide.shapes.add_textbox(
                Inches(0.8), y_pos,
                Inches(12), Inches(0.5)
            )
            tf = text_box.text_frame
            tf.text = suggestion
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme.TEXT
            p.word_wrap = True
            
            y_pos += Inches(0.6)
    
    def _add_final_slide(self):
        """添加结束页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme.PRIMARY
        bg.line.fill.background()
        
        text_box = slide.shapes.add_textbox(
            Inches(0), Inches(3),
            self.SLIDE_WIDTH, Inches(1)
        )
        tf = text_box.text_frame
        tf.text = "Thank You"
        p = tf.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.theme.WHITE
        p.alignment = PP_ALIGN.CENTER
        
        sub_box = slide.shapes.add_textbox(
            Inches(0), Inches(4.2),
            self.SLIDE_WIDTH, Inches(0.5)
        )
        tf = sub_box.text_frame
        tf.text = "Professional Layout Review Report"
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = self.theme.SECONDARY
        p.alignment = PP_ALIGN.CENTER
    
    def _add_slide_header(self, slide, title: str):
        """添加幻灯片标题栏"""
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(0.7)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme.PRIMARY
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.12),
            Inches(12), Inches(0.5)
        )
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.theme.WHITE
    
    def save(self, filepath: str):
        """保存PPT"""
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        self.prs.save(filepath)
        print(f"✓ PPTX report saved: {filepath}")


# ============================================================================
# PDF报告生成器
# ============================================================================

class PDFReportGenerator:
    """PDF格式报告生成器"""
    
    def __init__(self):
        if not REPORTLAB_AVAILABLE:
            raise ImportError("reportlab is required for PDF export")
        self.theme = ReportTheme()
        self.styles = getSampleStyleSheet()
        self._setup_styles()
    
    def _setup_styles(self):
        """设置PDF样式"""
        self.title_style = ParagraphStyle(
            'CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            textColor=self.theme.RL_PRIMARY,
            spaceAfter=30
        )
        
        self.heading_style = ParagraphStyle(
            'CustomHeading',
            parent=self.styles['Heading2'],
            fontSize=16,
            textColor=self.theme.RL_PRIMARY,
            spaceAfter=12
        )
        
        self.subheading_style = ParagraphStyle(
            'CustomSubHeading',
            parent=self.styles['Heading3'],
            fontSize=13,
            textColor=self.theme.RL_PRIMARY,
            spaceAfter=10
        )
        
        self.normal_style = ParagraphStyle(
            'CustomNormal',
            parent=self.styles['Normal'],
            fontSize=10,
            spaceAfter=6
        )
    
    def create_report(self, engine: ProfessionalLayoutReviewEngine,
                      title: str = "Layout Review Report") -> List:
        """创建PDF报告内容"""
        story = []
        summary = engine._generate_summary()
        
        # 封面
        story.extend(self._create_cover(title, engine.config))
        story.append(PageBreak())
        
        # 目录
        story.extend(self._create_toc())
        story.append(PageBreak())
        
        # Executive Summary
        story.extend(self._create_summary(summary, engine))
        story.append(PageBreak())
        
        # Net Statistics
        story.extend(self._create_net_statistics(engine))
        story.append(PageBreak())
        
        # Violations
        if summary.total_violations > 0:
            story.extend(self._create_violations_section(engine))
            story.append(PageBreak())
        
        # Matching Analysis
        if engine.matching_results:
            story.extend(self._create_matching_section(engine))
            story.append(PageBreak())
        
        # Recommendations
        story.extend(self._create_recommendations(engine))
        
        return story
    
    def _create_cover(self, title: str, config: LayoutReviewConfig) -> List:
        """创建封面"""
        story = []
        
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph(title, self.title_style))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(f"Technology: {config.tech_config.name}", self.heading_style))
        story.append(Paragraph(f"Process Node: {config.tech_config.node}", self.normal_style))
        story.append(Spacer(1, 1*inch))
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", self.normal_style))
        story.append(Paragraph("Professional Layout Review Engine v1.0", self.normal_style))
        
        return story
    
    def _create_toc(self) -> List:
        """创建目录"""
        story = []
        story.append(Paragraph("Table of Contents", self.heading_style))
        story.append(Spacer(1, 0.2*inch))
        
        sections = [
            "1. Executive Summary",
            "2. Net Statistics Overview",
            "3. Violation Analysis",
            "4. Matching Analysis",
            "5. Recommendations"
        ]
        
        for section in sections:
            story.append(Paragraph(section, self.normal_style))
        
        return story
    
    def _create_summary(self, summary: ReviewSummary, engine: ProfessionalLayoutReviewEngine) -> List:
        """创建Executive Summary"""
        story = []
        story.append(Paragraph("1. Executive Summary", self.heading_style))
        story.append(Spacer(1, 0.2*inch))
        
        # 统计表格
        data = [
            ['Metric', 'Value'],
            ['Total Nets', str(summary.total_nets)],
            ['Critical Issues', str(summary.critical_count)],
            ['Warnings', str(summary.warning_count)],
            ['Info Messages', str(summary.info_count)],
            ['Matching Pairs', str(summary.matching_pairs_analyzed)],
        ]
        
        table = Table(data, colWidths=[3*inch, 2*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), self.theme.RL_PRIMARY),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
        ]))
        
        story.append(table)
        story.append(Spacer(1, 0.3*inch))
        
        # RC Summary
        story.append(Paragraph("RC Analysis Summary", self.subheading_style))
        rc_data = [
            ['Parameter', 'Value'],
            ['Resistance Range', f"{summary.total_resistance_range[0]:.1f} - {summary.total_resistance_range[1]:.1f} Ω"],
            ['Average Resistance', f"{summary.avg_resistance:.1f} Ω"],
            ['Capacitance Range', f"{summary.total_capacitance_range[0]:.1f} - {summary.total_capacitance_range[1]:.1f} fF"],
            ['Average Capacitance', f"{summary.avg_capacitance:.1f} fF"],
        ]
        
        rc_table = Table(rc_data, colWidths=[3*inch, 2.5*inch])
        rc_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), self.theme.RL_PRIMARY),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
        ]))
        
        story.append(rc_table)
        
        return story
    
    def _create_net_statistics(self, engine: ProfessionalLayoutReviewEngine) -> List:
        """创建Net统计"""
        story = []
        story.append(Paragraph("2. Net Statistics Overview", self.heading_style))
        story.append(Spacer(1, 0.2*inch))
        
        table_data = engine.get_net_info_table()
        if not table_data:
            return story
        
        # 表头
        headers = ['Net Name', 'R (Ω)', 'C (fF)', 'Length', 'Vias', 'Violations']
        data = [headers]
        
        # 数据行 (最多20个)
        for row in table_data[:20]:
            data.append([
                row['Net Name'],
                f"{row['Total R (Ω)']:.2f}",
                f"{row['Total C (fF)']:.1f}",
                f"{row['Length (μm)']:.1f}",
                str(row['Via Count']),
                str(row['Violations'])
            ])
        
        table = Table(data, colWidths=[2.5*inch, 0.9*inch, 0.9*inch, 0.9*inch, 0.7*inch, 0.9*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), self.theme.RL_PRIMARY),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 9),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('ALIGN', (0, 0), (0, -1), 'LEFT'),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
        ]))
        
        story.append(table)
        
        return story
    
    def _create_violations_section(self, engine: ProfessionalLayoutReviewEngine) -> List:
        """创建违规分析"""
        story = []
        story.append(Paragraph("3. Violation Analysis", self.heading_style))
        story.append(Spacer(1, 0.2*inch))
        
        # 按严重程度分组
        critical = [v for v in engine.violations if v.severity == Severity.CRITICAL]
        warnings = [v for v in engine.violations if v.severity == Severity.WARNING]
        
        if critical:
            story.append(Paragraph("Critical Issues", self.subheading_style))
            for v in critical[:5]:  # 最多5个
                story.extend(self._create_violation_item(v))
            story.append(Spacer(1, 0.2*inch))
        
        if warnings:
            story.append(Paragraph("Warnings", self.subheading_style))
            for v in warnings[:5]:
                story.extend(self._create_violation_item(v))
        
        return story
    
    def _create_violation_item(self, v: Violation) -> List:
        """创建单个违规项"""
        story = []
        
        color = self.theme.RL_CRITICAL if v.severity == Severity.CRITICAL else \
                self.theme.RL_WARNING if v.severity == Severity.WARNING else self.theme.RL_INFO
        
        story.append(Paragraph(
            f"<font color='#{color.hexval()[2:8]}'><b>[{v.rule_id}]</b></font> {v.rule_name}",
            self.normal_style
        ))
        story.append(Paragraph(f"Net: <b>{v.net_name}</b>", self.normal_style))
        story.append(Paragraph(v.message, self.normal_style))
        if v.suggestion:
            story.append(Paragraph(f"<i>Suggestion: {v.suggestion}</i>", self.normal_style))
        story.append(Spacer(1, 0.1*inch))
        
        return story
    
    def _create_matching_section(self, engine: ProfessionalLayoutReviewEngine) -> List:
        """创建匹配分析"""
        story = []
        story.append(Paragraph("4. Matching Analysis", self.heading_style))
        story.append(Spacer(1, 0.2*inch))
        
        for match in engine.matching_results[:10]:  # 最多10个
            score_color = self.theme.RL_SUCCESS if match.match_score >= 80 else \
                         self.theme.RL_WARNING if match.match_score >= 60 else self.theme.RL_CRITICAL
            
            story.append(Paragraph(
                f"<b>{match.net1} ↔ {match.net2}</b> " +
                f"(<font color='#{score_color.hexval()[2:8]}'>Score: {match.match_score:.1f}</font>)",
                self.subheading_style
            ))
            
            data = [
                ['Metric', 'Value'],
                ['Length Ratio', f"{match.length_ratio:.3f}"],
                ['Resistance Ratio', f"{match.resistance_ratio:.3f}"],
                ['Via Count Diff', str(match.via_count_diff)],
            ]
            
            table = Table(data, colWidths=[2*inch, 1.5*inch])
            table.setStyle(TableStyle([
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
            ]))
            story.append(table)
            
            if match.issues:
                story.append(Paragraph("Issues:", self.normal_style))
                for issue in match.issues[:3]:
                    story.append(Paragraph(f"- {issue}", self.normal_style))
            
            story.append(Spacer(1, 0.2*inch))
        
        return story
    
    def _create_recommendations(self, engine: ProfessionalLayoutReviewEngine) -> List:
        """创建建议"""
        story = []
        story.append(Paragraph("5. Recommendations", self.heading_style))
        story.append(Spacer(1, 0.2*inch))
        
        suggestions = set()
        for v in engine.violations:
            if v.suggestion:
                suggestions.add(v.suggestion)
        
        for sug in sorted(suggestions)[:15]:
            story.append(Paragraph(f"• {sug}", self.normal_style))
        
        return story
    
    def save(self, story: List, filepath: str):
        """保存PDF"""
        doc = SimpleDocTemplate(
            filepath,
            pagesize=A4,
            rightMargin=0.75*inch,
            leftMargin=0.75*inch,
            topMargin=1*inch,
            bottomMargin=1*inch
        )
        doc.build(story)
        print(f"✓ PDF report saved: {filepath}")


# ============================================================================
# 便捷函数
# ============================================================================

def generate_reports(engine: ProfessionalLayoutReviewEngine,
                     output_dir: str,
                     base_name: str = "layout_review_report") -> Tuple[str, str]:
    """
    生成PPTX和PDF报告
    
    Returns:
        (pptx_path, pdf_path)
    """
    os.makedirs(output_dir, exist_ok=True)
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # PPTX报告
    pptx_path = os.path.join(output_dir, f"{base_name}_{timestamp}.pptx")
    pptx_gen = PPTXReportGenerator()
    pptx_gen.create_report(engine)
    pptx_gen.save(pptx_path)
    
    # PDF报告
    pdf_path = os.path.join(output_dir, f"{base_name}_{timestamp}.pdf")
    if REPORTLAB_AVAILABLE:
        pdf_gen = PDFReportGenerator()
        story = pdf_gen.create_report(engine)
        pdf_gen.save(story, pdf_path)
    else:
        print("Warning: PDF generation requires reportlab")
        pdf_path = None
    
    return pptx_path, pdf_path


if __name__ == '__main__':
    print("Professional Report Generator")
    print("=" * 60)
    print("\nAvailable formats:")
    print("  - PPTX: Microsoft PowerPoint")
    print(f"  - PDF:  Adobe PDF {'(available)' if REPORTLAB_AVAILABLE else '(requires reportlab)'}")


def generate_routing_report(state, app_state, output_dir="./output", base_name="routing_report"):
    """Thin wrapper for the routing PPTX report (compatibility shim)."""
    from report.routing_pptx import generate_routing_pptx
    import os
    from datetime import datetime
    os.makedirs(output_dir, exist_ok=True)
    path = os.path.join(output_dir, f"{base_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
    return generate_routing_pptx(state, app_state, path)
