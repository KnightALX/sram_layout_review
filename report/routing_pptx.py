"""Routing-review PPTX report.

Generates a slide deck with:
- Slide 1: Cover (title, preset, timestamp)
- Slide 2: Executive Summary (6 metric averages)
- Slide 3: Golden Net detail (table + image placeholder)
- Slide 4+: Per-batch-net page (3-section layout: full / directional / violation)
- Final: Recommendations
"""
from __future__ import annotations
import os
from datetime import datetime
from typing import TYPE_CHECKING
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

if TYPE_CHECKING:
    from app.routing_state import RoutingState
    from app.state import AppState


def _add_title_slide(prs, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_summary_slide(prs, state: "RoutingState"):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Executive Summary"
    n = len(state.batch_results) or 1
    avg_h = sum(r["h_ratio"] for r in state.batch_results.values()) / n
    avg_v = sum(r["v_ratio"] for r in state.batch_results.values()) / n
    avg_r = sum(r["r_total"] for r in state.batch_results.values()) / n
    avg_tau = sum(r["effective_tau_ps"] for r in state.batch_results.values()) / n
    avg_sim = sum(r["similarity_score"] for r in state.batch_results.values()) / n
    pass_pct = sum(1 for r in state.batch_results.values() if r["gate_pass"]) / n * 100
    missing_total = sum(r["missing_via_count"] for r in state.batch_results.values())
    rows = [
        ("Avg H Ratio", f"{avg_h*100:.1f}%"),
        ("Avg V Ratio", f"{avg_v*100:.1f}%"),
        ("Avg Eff. R", f"{avg_r:.2f} Ω"),
        ("Avg Eff. τ", f"{avg_tau:.2f} ps"),
        ("Avg Similarity", f"{avg_sim:.1f}/100"),
        ("Pass Rate", f"{pass_pct:.0f}%"),
        ("Total Missing Vias", str(missing_total)),
    ]
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = body.text_frame
    for i, (k, v) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{k}: {v}"
        p.font.size = Pt(18)


def _add_golden_slide(prs, state: "RoutingState"):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Golden Net: {state.golden_net_name or '(none)'}"
    if not state.golden_metrics:
        return
    m = state.golden_metrics
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = body.text_frame
    tf.text = "Feature Vector (Golden Reference):"
    for k, v in m.items():
        p = tf.add_paragraph()
        p.text = f"  {k}: {v}"
        p.font.size = Pt(16)


def _add_net_slide(prs, name: str, m: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Net: {name}"
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = body.text_frame
    tf.text = f"Pass: {'YES' if m['gate_pass'] else 'NO'}"
    lines = [
        f"H length: {m['h_len']:.2f} μm ({m['h_ratio']*100:.1f}%)",
        f"V length: {m['v_len']:.2f} μm ({m['v_ratio']*100:.1f}%)",
        f"Dominant: {m['dominant']}",
        f"Total R: {m['r_total']:.2f} Ω",
        f"Total C: {m['c_total']:.2f} fF",
        f"Effective τ: {m['effective_tau_ps']:.2f} ps",
        f"Via coverage: {m['via_coverage']*100:.1f}%",
        f"Missing vias: {m['missing_via_count']}",
        f"Similarity to Golden: {m['similarity_score']:.1f}/100",
    ]
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
    if m["gate_fail_reasons"]:
        p = tf.add_paragraph()
        p.text = "Fail reasons:"
        p.font.size = Pt(14)
        p.font.bold = True
        for r in m["gate_fail_reasons"]:
            p = tf.add_paragraph()
            p.text = f"  - {r}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)


def generate_routing_pptx(
    state: "RoutingState",
    app_state: "AppState",
    out_path: str,
):
    """Generate the routing review PPTX report."""
    prs = Presentation()
    _add_title_slide(
        prs,
        "Routing Review Report",
        f"Preset: {state.current_preset} | "
        f"{len(state.batch_results)} nets | "
        f"{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
    )
    _add_summary_slide(prs, state)
    if state.golden_net_name:
        _add_golden_slide(prs, state)
    for name, m in state.batch_results.items():
        _add_net_slide(prs, name, m)
    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    return out_path
